from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm
from pptx.util import Pt

class GeneratePPTX():    
    color_map = {"blue":RGBColor(0, 112, 192), 
            "orange":RGBColor(237, 125, 49), 
            "green":RGBColor(112, 173, 71), 
            "red":RGBColor(255, 0, 0),
            'black':RGBColor(0, 0, 0), 
            'white':RGBColor(255, 255, 255), 
            'yellow':RGBColor(255, 192, 0),
           }
    
    def __init__(self, path=None):
        self.path = path
        self.prs = Presentation(path)
    
    def set_slide_size(self, slide_width, slide_height):
        self.prs.slide_width = Cm(slide_width)
        self.prs.slide_height = Cm(slide_height)
    
    def num_slides(self):
        '''
        슬라이드의 개수를 반환한다.
        '''
        return len(self.prs.slides)
    
    def num_shapes(self, slide):
        '''
        슬라이드 내 도형의 개수를 반환한다.
        
        input
        ----
        slide: 슬라이드 번호        
        '''
        return len(self.prs.slides[slide-1].shapes)
    
    def shape_name(self, slide):
        '''
        슬라이드 내 도형의 {이름:번호} 사전을 반환한다.
        
        input
        ----
        slide: 슬라이드 번호
        '''
        return {s.name:i for i, s in enumerate(self.prs.slides[slide-1].shapes)}
    
    def add_slide(self, layout=6):
        '''
        슬라이드를 새로 추가한다.
        ------
        layout: 추가되는 슬라이드의 레이아웃을 결정한다. 다음의 0~11 중 선택. 
        
                제목 슬라이드(0), 제목 및 내용(1), 구역 머리글(2), 콘텐츠 2개(3), 
                비교(4), 제목만(5), 빈 화면(6), 캡션 있는 콘텐츠(7), 캡션 있는 그림(8), 
                제목 및 세로 텍스트(9), 세로 제목 및 텍스트(10)
        '''
        self.prs.slides.add_slide(self.prs.slide_layouts[layout])
                 
    def add_picture(self, slide, img_path, left, top, width=None, height=None):
        '''
        슬라이드에 이미지를 추가한다.
        ------
        slide: 슬라이드 번호
        img_path: 이미지 경로
        left: 이미지와 슬라이드 왼쪽 테두리 사이의 간격(cm)
        top: 이미지와 슬라이드 위쪽 테두리 사이의 간격(cm)
        width: 이미지 너비
        height: 이미지 높이
        '''
        slide = self.prs.slides[slide-1]
        
        width = None if width is None else Cm(width)
        height = None if height is None else Cm(height)
        picture = slide.shapes.add_picture(img_path, Cm(left), Cm(top),
                                               width, height)
    
    def add_textbox(self, slide, left, top, width, height, name=None):
        slide = self.prs.slides[slide-1]
        txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
        if name is not None:
            txBox.name = name
            
    def add_table(self, slide, rows, cols, left, top, width, height, name=None):
        slide = self.prs.slides[slide-1]
        table = slide.shapes.add_table(rows, cols, Cm(left), Cm(top), Cm(width), Cm(height))
        if name is not None:
            table.name = name
            
    def merge_table_cells(self, slide, shape, cell_origin, cell_spanned):    
        cell0 = self._get_shape(slide, shape, cell_origin)
        cell1 = self._get_shape(slide, shape, cell_spanned)
        cell0.merge(cell1)
    
    def show_text(self, slide, shape, cell=None):

        shape = self._get_shape(slide, shape, cell=cell)

        tf = shape.text_frame
        return tf.text
        
    def edit_text(self, slide, shape, text, cell=None, clear=True, hlink=None,
                 bold=False, italic=False, underline=False, size=32, font=None, color='black',
                 level=0, alignment='left'):
        '''
        도형에 텍스트를 추가하거나 삭제한다.
        
        input
        ----
        slide: 슬라이드 번호(int)
        shape: 도형 이름(str), 또는 도형 번호(int)
        text: 써넣을 텍스트 (str)
        cell: 도형의 종류가 표인 경우, 표의 어느 부분에 글씨를 써 넣을지 지정할 수 있다.
            예를 들어 표 (0, 1)번 셀에 글씨를 추가하고 싶은 경우 cell=(0, 1)로 지정한다.
        clear: 기존에 있던 글씨를 지울지, 또는 이어서 글을 추가할지 여부를 결정한다.
        hlink: 하이퍼 링크 주소(str)
        bold: 볼드체 여부
        italic: 이탤릭체 여부
        underline: 밑줄 여부
        size: 글씨 크기
        font: 폰트 이름
        color: 글씨 색상 이름(str), 또는 RGB 색상(tuple)
        level: 단락의 단위. (level이 1 증가할 때마다 들여쓰기가 시행된다)
        alignment: 정렬 기준. (left, center, right)
        '''
        
        shape = self._get_shape(slide, shape, cell=cell)
            
        tf = shape.text_frame
        
        if clear:
            tf.clear()
        
        p = tf.paragraphs[0]
        p.level = level
        if alignment == 'left':
            p.alignment = PP_ALIGN.LEFT
        elif alignment == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif alignment == 'right':
            p.alignment = PP_ALIGN.RIGHT
        
        r = p.add_run()
        r.text = text
        r.font.bold = bold
        r.font.italic = italic
        r.font.underline = underline
        r.font.size = Pt(size)
        r.font.name = font
        r.font.color.rgb = self.cmap(color)
        r.hyperlink.address = hlink
    
    def change_bgcolor(self, slide, shape, bg_color, cell=None):
        '''
        도형의 배경 색을 바꾼다. 
        
        input
        ----
        slide: 슬라이드 번호(int)
        shape: 도형 이름(str), 또는 도형 번호(int)
        bg_color: 색상 이름(str), 또는 RGB 색상(tuple)
        cell: 도형의 종류가 표인 경우, 표의 어느 부분에 글씨를 써 넣을지 지정할 수 있다.
            예를 들어 표 (0, 1)번 셀에 글씨를 추가하고 싶은 경우 cell=(0, 1)로 지정한다.
        '''
        shape = self._get_shape(slide, shape, cell=cell)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = self.cmap(bg_color)  

    def _get_shape(self, slide, shape, cell=None):    
        if isinstance(shape, str):
            shape_name_to_id = self.shape_name(slide)
            slide = self.prs.slides[slide-1]
            shape = slide.shapes[shape_name_to_id[shape]]

        elif isinstance(shape, int):
            slide = self.prs.slides[slide-1]
            shape = slide.shapes[shape]
            
        if cell is not None:
            shape = shape.table.cell(*cell)
            
        return shape
            
    @classmethod
    def cmap(cls, color):
        
        if isinstance(color, str):
            color = cls.color_map[color.lower()]
            
        elif isinstance(color, tuple):
            color = RGBColor(*color)
        
        return color

    def save(self, savepath):
        self.prs.save(savepath)
